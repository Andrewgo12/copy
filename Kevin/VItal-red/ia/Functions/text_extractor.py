"""
Text Extraction Engine
Extracts text from various file formats and email bodies
"""

import email
import os
import re
import datetime
from typing import Dict, List, Any, Optional
from email.message import Message
import logging
import mimetypes

logger = logging.getLogger(__name__)

class TextExtractor:
    """
    Extracts text from various sources including email bodies and attachments
    """
    
    def __init__(self, base_path: str):
        """
        Initialize text extractor
        
        Args:
            base_path: Base path for the ia folder
        """
        self.base_path = base_path
        self.text_path = os.path.join(base_path, "Text")
        os.makedirs(self.text_path, exist_ok=True)
    
    def extract_email_body(self, email_message: Message) -> Dict[str, str]:
        """
        Extract text and HTML content from email body
        
        Args:
            email_message: Email message object
            
        Returns:
            Dict: Contains 'text' and 'html' content
        """
        body_content = {
            'text': '',
            'html': '',
            'raw': ''
        }
        
        try:
            if email_message.is_multipart():
                # Handle multipart messages
                for part in email_message.walk():
                    content_type = part.get_content_type()
                    content_disposition = part.get_content_disposition()
                    
                    # Skip attachments
                    if content_disposition == 'attachment':
                        continue
                    
                    if content_type == "text/plain":
                        payload = part.get_payload(decode=True)
                        if isinstance(payload, bytes):
                            text_content = payload.decode('utf-8', errors='ignore')
                        else:
                            text_content = str(payload)
                        body_content['text'] += text_content + '\n'
                    
                    elif content_type == "text/html":
                        payload = part.get_payload(decode=True)
                        if isinstance(payload, bytes):
                            html_content = payload.decode('utf-8', errors='ignore')
                        else:
                            html_content = str(payload)
                        body_content['html'] += html_content + '\n'
            
            else:
                # Handle single part messages
                content_type = email_message.get_content_type()
                payload = email_message.get_payload(decode=True)
                
                if isinstance(payload, bytes):
                    content = payload.decode('utf-8', errors='ignore')
                else:
                    content = str(payload)
                
                if content_type == "text/plain":
                    body_content['text'] = content
                elif content_type == "text/html":
                    body_content['html'] = content
                else:
                    body_content['raw'] = content
            
            # Clean up text content
            body_content['text'] = self.clean_text(body_content['text'])
            
            # Extract text from HTML if no plain text available
            if not body_content['text'] and body_content['html']:
                body_content['text'] = self.html_to_text(body_content['html'])
            
        except Exception as e:
            logger.error(f"Error extracting email body: {str(e)}")
        
        return body_content
    
    def clean_text(self, text: str) -> str:
        """
        Clean and normalize text content with professional formatting

        Args:
            text: Raw text content

        Returns:
            str: Professionally cleaned and formatted text
        """
        if not text:
            return ""

        try:
            # Step 1: Normalize line endings
            text = text.replace('\r\n', '\n').replace('\r', '\n')

            # Step 2: Remove excessive whitespace while preserving structure
            lines = text.split('\n')
            cleaned_lines = []

            for line in lines:
                # Clean each line individually
                cleaned_line = re.sub(r'\s+', ' ', line).strip()

                # Skip quoted lines (email replies) but preserve important quotes
                if cleaned_line.startswith('>'):
                    # Only skip if it's clearly a reply quote, not a legitimate quote
                    if not any(keyword in cleaned_line.lower() for keyword in ['important', 'note', 'attention']):
                        continue

                # Skip lines that are just separators
                if re.match(r'^[-=_*]{3,}$', cleaned_line):
                    continue

                # Skip forwarded email headers
                if re.match(r'^\s*(From|Sent|To|Subject|Date):\s*', cleaned_line, re.IGNORECASE):
                    continue

                cleaned_lines.append(cleaned_line)

            # Step 3: Rejoin and format professionally
            cleaned_text = '\n'.join(cleaned_lines)

            # Step 4: Normalize paragraph spacing
            cleaned_text = re.sub(r'\n\s*\n\s*\n+', '\n\n', cleaned_text)  # Multiple newlines to double

            # Step 5: Clean up common email artifacts
            cleaned_text = re.sub(r'Sent from my \w+', '', cleaned_text, flags=re.IGNORECASE)
            cleaned_text = re.sub(r'Get Outlook for \w+', '', cleaned_text, flags=re.IGNORECASE)
            cleaned_text = re.sub(r'Confidentiality Notice:.*$', '', cleaned_text, flags=re.IGNORECASE | re.DOTALL)

            # Step 6: Preserve important formatting markers
            cleaned_text = self._preserve_formatting_markers(cleaned_text)

            return cleaned_text.strip()

        except Exception as e:
            logger.warning(f"Error cleaning text: {e}")
            # Fallback to basic cleaning
            return re.sub(r'\s+', ' ', text).strip()

    def _preserve_formatting_markers(self, text: str) -> str:
        """Preserve important formatting markers in text"""
        try:
            # Preserve bullet points and lists
            text = re.sub(r'^(\s*)[-*•]\s+', r'\1• ', text, flags=re.MULTILINE)

            # Preserve numbered lists
            text = re.sub(r'^(\s*)(\d+)\.?\s+', r'\1\2. ', text, flags=re.MULTILINE)

            # Preserve section headers (lines that are all caps or have colons)
            lines = text.split('\n')
            formatted_lines = []

            for line in lines:
                stripped = line.strip()
                if stripped:
                    # Check if it's a header (all caps, ends with colon, or is short and prominent)
                    if (stripped.isupper() and len(stripped) < 50 and len(stripped) > 3) or stripped.endswith(':'):
                        formatted_lines.append(f"\n{stripped}")
                        formatted_lines.append("-" * min(len(stripped), 40))
                    else:
                        formatted_lines.append(line)
                else:
                    formatted_lines.append(line)

            return '\n'.join(formatted_lines)

        except Exception as e:
            logger.warning(f"Error preserving formatting markers: {e}")
            return text
    
    def html_to_text(self, html: str) -> str:
        """
        Convert HTML to plain text (basic implementation)
        
        Args:
            html: HTML content
            
        Returns:
            str: Plain text content
        """
        if not html:
            return ""
        
        try:
            # Try to use html2text if available
            try:
                import html2text
                h = html2text.HTML2Text()
                h.ignore_links = True
                h.ignore_images = True
                return h.handle(html)
            except ImportError:
                pass
            
            # Fallback: basic HTML tag removal
            import html
            
            # Decode HTML entities
            text = html.unescape(html)
            
            # Remove script and style elements
            text = re.sub(r'<script[^>]*>.*?</script>', '', text, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
            
            # Convert common HTML elements to text
            text = re.sub(r'<br[^>]*>', '\n', text, flags=re.IGNORECASE)
            text = re.sub(r'<p[^>]*>', '\n', text, flags=re.IGNORECASE)
            text = re.sub(r'</p>', '\n', text, flags=re.IGNORECASE)
            text = re.sub(r'<div[^>]*>', '\n', text, flags=re.IGNORECASE)
            text = re.sub(r'</div>', '\n', text, flags=re.IGNORECASE)
            
            # Remove all remaining HTML tags
            text = re.sub(r'<[^>]+>', '', text)
            
            # Clean up whitespace
            text = self.clean_text(text)
            
            return text
            
        except Exception as e:
            logger.error(f"Error converting HTML to text: {str(e)}")
            return ""
    
    def extract_from_pdf(self, file_path: str) -> str:
        """
        Extract text from PDF file
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            str: Extracted text
        """
        try:
            # Try PyPDF2 first
            try:
                import PyPDF2
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                    return self.clean_text(text)
            except ImportError:
                pass
            
            # Try pdfplumber as alternative
            try:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    text = ""
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                    return self.clean_text(text)
            except ImportError:
                pass
            
            logger.warning(f"No PDF extraction library available for {file_path}")
            return ""
            
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {str(e)}")
            return ""
    
    def extract_from_docx(self, file_path: str) -> str:
        """
        Extract text from DOCX file
        
        Args:
            file_path: Path to DOCX file
            
        Returns:
            str: Extracted text
        """
        try:
            import docx
            doc = docx.Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return self.clean_text(text)
            
        except ImportError:
            logger.warning(f"python-docx not available for {file_path}")
            return ""
        except Exception as e:
            logger.error(f"Error extracting text from DOCX {file_path}: {str(e)}")
            return ""
    
    def extract_from_xlsx(self, file_path: str) -> str:
        """
        Extract text from XLSX file
        
        Args:
            file_path: Path to XLSX file
            
        Returns:
            str: Extracted text
        """
        try:
            import openpyxl
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        text += "\t".join(row_text) + "\n"
                text += "\n"
            
            return self.clean_text(text)
            
        except ImportError:
            logger.warning(f"openpyxl not available for {file_path}")
            return ""
        except Exception as e:
            logger.error(f"Error extracting text from XLSX {file_path}: {str(e)}")
            return ""
    
    def extract_from_image(self, file_path: str) -> str:
        """
        Extract text from image using OCR
        
        Args:
            file_path: Path to image file
            
        Returns:
            str: Extracted text
        """
        try:
            import pytesseract
            from PIL import Image
            
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return self.clean_text(text)
            
        except ImportError:
            logger.warning(f"OCR libraries not available for {file_path}")
            return ""
        except Exception as e:
            logger.error(f"Error extracting text from image {file_path}: {str(e)}")
            return ""
    
    def extract_from_file(self, file_path: str) -> str:
        """
        Extract text from file based on its type
        
        Args:
            file_path: Path to file
            
        Returns:
            str: Extracted text
        """
        if not os.path.exists(file_path):
            return ""
        
        # Get file extension
        _, ext = os.path.splitext(file_path.lower())
        
        # Extract based on file type
        if ext == '.pdf':
            return self.extract_from_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            return self.extract_from_docx(file_path)
        elif ext in ['.xlsx', '.xls']:
            return self.extract_from_xlsx(file_path)
        elif ext in ['.txt', '.log', '.csv']:
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return self.clean_text(f.read())
            except Exception as e:
                logger.error(f"Error reading text file {file_path}: {str(e)}")
                return ""
        elif ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif']:
            return self.extract_from_image(file_path)
        else:
            logger.info(f"Unsupported file type for text extraction: {ext}")
            return ""
    
    def save_extracted_text(self, unique_id: str, content: Dict[str, Any]) -> str:
        """
        Save extracted text content to professionally formatted file

        Args:
            unique_id: Unique email identifier
            content: Extracted content dictionary

        Returns:
            str: Path to saved text file
        """
        email_text_folder = os.path.join(self.text_path, unique_id)
        os.makedirs(email_text_folder, exist_ok=True)

        text_file_path = os.path.join(email_text_folder, "extracted_text.txt")

        try:
            with open(text_file_path, 'w', encoding='utf-8') as f:
                # Professional header
                f.write("EMAIL COMMUNICATION CONTENT EXTRACTION\n")
                f.write("=" * 70 + "\n")
                f.write(f"Document ID: {unique_id}\n")
                f.write(f"Extraction Date: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
                f.write(f"Processing System: Gmail Data Extraction System v2.0\n")
                f.write("=" * 70 + "\n\n")

                # Table of Contents
                f.write("TABLE OF CONTENTS:\n")
                f.write("-" * 30 + "\n")
                f.write("1. Email Body Content\n")
                if content.get('attachments'):
                    f.write("2. Attachment Content\n")
                    for i, attachment in enumerate(content['attachments'], 1):
                        if attachment.get('text'):
                            f.write(f"   2.{i} {attachment.get('filename', 'Unknown File')}\n")
                f.write("3. Content Summary\n")
                f.write("\n" + "=" * 70 + "\n\n")

                # Section 1: Email Body Content
                f.write("1. EMAIL BODY CONTENT\n")
                f.write("=" * 30 + "\n\n")

                email_body = content.get('email_body', {})
                if email_body.get('text'):
                    cleaned_text = self.clean_text(email_body['text'])
                    f.write("Plain Text Content:\n")
                    f.write("-" * 20 + "\n")
                    f.write(cleaned_text)
                    f.write("\n\n")

                    # Word count and statistics
                    word_count = len(cleaned_text.split())
                    char_count = len(cleaned_text)
                    f.write(f"Content Statistics:\n")
                    f.write(f"- Word Count: {word_count:,}\n")
                    f.write(f"- Character Count: {char_count:,}\n")
                    f.write(f"- Estimated Reading Time: {max(1, word_count // 200)} minutes\n\n")
                else:
                    f.write("No plain text content available.\n\n")

                if email_body.get('html'):
                    f.write("HTML Content (converted to text):\n")
                    f.write("-" * 35 + "\n")
                    html_text = self.html_to_text(email_body['html'])
                    cleaned_html_text = self.clean_text(html_text)
                    f.write(cleaned_html_text)
                    f.write("\n\n")

                # Section 2: Attachment Content
                if content.get('attachments'):
                    f.write("2. ATTACHMENT CONTENT\n")
                    f.write("=" * 30 + "\n\n")

                    for i, attachment in enumerate(content['attachments'], 1):
                        filename = attachment.get('filename', 'Unknown File')
                        extracted_text = attachment.get('text', '')

                        f.write(f"2.{i} ATTACHMENT: {filename}\n")
                        f.write("-" * (15 + len(filename)) + "\n")

                        if extracted_text:
                            cleaned_attachment_text = self.clean_text(extracted_text)
                            f.write(f"Content:\n")
                            f.write(cleaned_attachment_text)
                            f.write("\n\n")

                            # Attachment statistics
                            att_word_count = len(cleaned_attachment_text.split())
                            att_char_count = len(cleaned_attachment_text)
                            f.write(f"Attachment Statistics:\n")
                            f.write(f"- Word Count: {att_word_count:,}\n")
                            f.write(f"- Character Count: {att_char_count:,}\n")
                            f.write(f"- File Type: {os.path.splitext(filename)[1].upper() or 'Unknown'}\n\n")
                        else:
                            f.write("No extractable text content found in this attachment.\n\n")

                        f.write("-" * 70 + "\n\n")

                # Section 3: Content Summary
                f.write("3. CONTENT SUMMARY\n")
                f.write("=" * 30 + "\n\n")

                total_words = 0
                total_chars = 0

                # Count email body
                if email_body.get('text'):
                    email_words = len(self.clean_text(email_body['text']).split())
                    email_chars = len(self.clean_text(email_body['text']))
                    total_words += email_words
                    total_chars += email_chars

                # Count attachments
                attachment_count = 0
                successful_extractions = 0
                for attachment in content.get('attachments', []):
                    attachment_count += 1
                    if attachment.get('text'):
                        successful_extractions += 1
                        att_text = self.clean_text(attachment['text'])
                        total_words += len(att_text.split())
                        total_chars += len(att_text)

                f.write(f"Overall Statistics:\n")
                f.write(f"- Total Word Count: {total_words:,}\n")
                f.write(f"- Total Character Count: {total_chars:,}\n")
                f.write(f"- Total Attachments: {attachment_count}\n")
                f.write(f"- Successful Text Extractions: {successful_extractions}\n")
                f.write(f"- Extraction Success Rate: {(successful_extractions/max(1,attachment_count)*100):.1f}%\n")
                f.write(f"- Estimated Total Reading Time: {max(1, total_words // 200)} minutes\n\n")

                f.write("Content Quality Assessment:\n")
                if total_words > 100:
                    f.write("- Content Volume: Substantial\n")
                elif total_words > 20:
                    f.write("- Content Volume: Moderate\n")
                else:
                    f.write("- Content Volume: Minimal\n")

                if successful_extractions == attachment_count and attachment_count > 0:
                    f.write("- Extraction Quality: Excellent (100% success)\n")
                elif successful_extractions > 0:
                    f.write("- Extraction Quality: Good (partial success)\n")
                else:
                    f.write("- Extraction Quality: Limited (text only)\n")

                f.write("\n" + "=" * 70 + "\n")
                f.write("End of Document\n")

            return text_file_path

        except Exception as e:
            logger.error(f"Error saving extracted text: {str(e)}")
            return ""

    def extract_from_pptx(self, file_path: str) -> str:
        """
        Extract text from PowerPoint PPTX file

        Args:
            file_path: Path to PPTX file

        Returns:
            str: Extracted text
        """
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""

            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"

            return self.clean_text(text)

        except ImportError:
            logger.warning(f"python-pptx not available for {file_path}")
            return ""
        except Exception as e:
            logger.error(f"Error extracting text from PPTX {file_path}: {str(e)}")
            return ""

    def extract_from_rtf(self, file_path: str) -> str:
        """
        Extract text from RTF file

        Args:
            file_path: Path to RTF file

        Returns:
            str: Extracted text
        """
        try:
            from striprtf.striprtf import rtf_to_text
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            text = rtf_to_text(rtf_content)
            return self.clean_text(text)

        except ImportError:
            logger.warning(f"striprtf not available for {file_path}")
            # Fallback: basic RTF parsing
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                # Very basic RTF text extraction
                import re
                text = re.sub(r'\\[a-z]+\d*\s?', '', content)  # Remove RTF commands
                text = re.sub(r'[{}]', '', text)  # Remove braces
                return self.clean_text(text)
            except Exception as e:
                logger.error(f"Error with fallback RTF extraction: {str(e)}")
                return ""
        except Exception as e:
            logger.error(f"Error extracting text from RTF {file_path}: {str(e)}")
            return ""

    def extract_from_csv(self, file_path: str) -> str:
        """
        Extract text from CSV file with proper formatting

        Args:
            file_path: Path to CSV file

        Returns:
            str: Extracted text
        """
        try:
            import csv
            text = ""

            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                # Try to detect delimiter
                sample = f.read(1024)
                f.seek(0)
                sniffer = csv.Sniffer()
                delimiter = sniffer.sniff(sample).delimiter

                reader = csv.reader(f, delimiter=delimiter)
                for row_num, row in enumerate(reader, 1):
                    if row_num == 1:
                        text += "Headers: " + " | ".join(row) + "\n"
                        text += "-" * 50 + "\n"
                    else:
                        text += f"Row {row_num-1}: " + " | ".join(row) + "\n"

            return self.clean_text(text)

        except Exception as e:
            logger.error(f"Error extracting text from CSV {file_path}: {str(e)}")
            return ""

    def extract_with_language_detection(self, text: str) -> Dict[str, Any]:
        """
        Extract text with language detection and analysis

        Args:
            text: Input text

        Returns:
            Dict: Text analysis results
        """
        analysis = {
            'text': text,
            'language': 'unknown',
            'confidence': 0.0,
            'word_count': 0,
            'character_count': 0,
            'sentence_count': 0,
            'readability_score': 0.0,
            'encoding': 'utf-8'
        }

        if not text:
            return analysis

        try:
            # Language detection
            try:
                from langdetect import detect, detect_langs
                analysis['language'] = detect(text)
                lang_probs = detect_langs(text)
                if lang_probs:
                    analysis['confidence'] = lang_probs[0].prob
            except ImportError:
                logger.warning("langdetect not available for language detection")
            except:
                pass  # Language detection failed

            # Basic text statistics
            analysis['word_count'] = len(text.split())
            analysis['character_count'] = len(text)

            # Sentence count (basic)
            import re
            sentences = re.split(r'[.!?]+', text)
            analysis['sentence_count'] = len([s for s in sentences if s.strip()])

            # Readability score (basic implementation)
            try:
                import textstat
                analysis['readability_score'] = textstat.flesch_reading_ease(text)
            except ImportError:
                # Basic readability calculation
                if analysis['sentence_count'] > 0 and analysis['word_count'] > 0:
                    avg_sentence_length = analysis['word_count'] / analysis['sentence_count']
                    # Simplified readability score
                    analysis['readability_score'] = max(0, 100 - avg_sentence_length * 2)

        except Exception as e:
            logger.warning(f"Error in text analysis: {str(e)}")

        return analysis

    def extract_structured_content(self, file_path: str) -> Dict[str, Any]:
        """
        Extract structured content with metadata and analysis

        Args:
            file_path: Path to file

        Returns:
            Dict: Structured content with analysis
        """
        result = {
            'file_path': file_path,
            'file_type': '',
            'extraction_method': '',
            'raw_text': '',
            'analysis': {},
            'extraction_successful': False,
            'extraction_time': 0,
            'error_message': ''
        }

        if not os.path.exists(file_path):
            result['error_message'] = "File not found"
            return result

        import time
        start_time = time.time()

        try:
            # Determine file type and extract
            _, ext = os.path.splitext(file_path.lower())
            result['file_type'] = ext

            if ext == '.pdf':
                result['extraction_method'] = 'pdf_extraction'
                result['raw_text'] = self.extract_from_pdf(file_path)
            elif ext in ['.docx', '.doc']:
                result['extraction_method'] = 'docx_extraction'
                result['raw_text'] = self.extract_from_docx(file_path)
            elif ext in ['.xlsx', '.xls']:
                result['extraction_method'] = 'xlsx_extraction'
                result['raw_text'] = self.extract_from_xlsx(file_path)
            elif ext == '.pptx':
                result['extraction_method'] = 'pptx_extraction'
                result['raw_text'] = self.extract_from_pptx(file_path)
            elif ext == '.rtf':
                result['extraction_method'] = 'rtf_extraction'
                result['raw_text'] = self.extract_from_rtf(file_path)
            elif ext == '.csv':
                result['extraction_method'] = 'csv_extraction'
                result['raw_text'] = self.extract_from_csv(file_path)
            elif ext in ['.txt', '.log']:
                result['extraction_method'] = 'text_file'
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    result['raw_text'] = self.clean_text(f.read())
            elif ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif']:
                result['extraction_method'] = 'ocr_extraction'
                result['raw_text'] = self.extract_from_image(file_path)
            else:
                result['error_message'] = f"Unsupported file type: {ext}"
                return result

            # Perform text analysis
            if result['raw_text']:
                result['analysis'] = self.extract_with_language_detection(result['raw_text'])
                result['extraction_successful'] = True
            else:
                result['error_message'] = "No text extracted"

        except Exception as e:
            result['error_message'] = str(e)
            logger.error(f"Error in structured content extraction: {str(e)}")

        finally:
            result['extraction_time'] = time.time() - start_time

        return result
